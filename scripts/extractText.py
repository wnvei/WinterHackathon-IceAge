import os 
import sys 
from pathlib import Path
from typing import List,Tuple

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try: 
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_text_from_pdf(file_path: str) -> str:
    if PyPDF2 is None:
          return f"PyPDF2 library is not installed. Cannot extract text from PDF."
    try:
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {page_num} ---\n{page_text}\n")
        return "\n".join(text) if text else "[No text extracted from PDF]\n"
    except Exception as e:
        return f"[ERROR extracting PDF {file_path}: {str(e)}]\n"
    
def extract_text_from_docx(file_path: str) -> str:
    if Document is None:
        return f"python-docx library is not installed. Cannot extract text from DOCX."
    try:
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text) if text else "[No text extracted from DOCX]\n"
    except Exception as e:
        return f"[ERROR extracting DOCX {file_path}: {str(e)}]\n"
    
def extract_text_from_pptx(file_path: str) -> str:
    if Presentation is None:
        return f"python-pptx library is not installed. Cannot extract text from PPTX."
    try:
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text) + "\n")
        return "\n".join(text) if text else "[No text extracted from PPTX]\n"
    except Exception as e:
        return f"[ERROR extracting PPTX {file_path}: {str(e)}]\n"
    
def extract_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        return f"[Unsupported file type: {ext}]\n"

def walk_and_extract(source_dir: str, output_file: str) -> Tuple[int, int]:
    files_processed = 0
    files_succeeded = 0
    all_text = []
    supported_extensions = {'.pdf', '.docx', '.pptx'}
    
    if not os.path.exists(source_dir):
        print(f"Warning: Directory {source_dir} does not exist")
        return 0, 0
    
    for root, dirs, files in os.walk(source_dir):
        for file in files:
            file_path = os.path.join(root, file)
            ext = Path(file).suffix.lower()
            
            if ext in supported_extensions:
                files_processed += 1
                print(f"Processing: {file_path}")
                
                relative_path = os.path.relpath(file_path, source_dir)
                all_text.append(f"\n{'='*80}")
                all_text.append(f"FILE: {relative_path}")
                all_text.append(f"{'='*80}\n")
                
                extracted_text = extract_text_from_file(file_path)
                all_text.append(extracted_text)
                if not extracted_text.startswith("[ERROR"):
                    files_succeeded += 1

    if all_text:
        os.makedirs(os.path.dirname(output_file), exist_ok=True)
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(all_text))
        print(f"Saved to: {output_file}")

    return files_processed, files_succeeded

def main():
    base_dir = Path(__file__).parent.parent 
    data_dir = base_dir / "Subject"

    #Write a pipeline to combine all the scripts and extract target subject argument later
    target_subject = sys.argv[1] if len(sys.argv) > 1 else None

    if not data_dir.exists():
        print(f"Data directory {data_dir} does not exist.")
        return
    
    for subject_dir in data_dir.iterdir():
        if not subject_dir.is_dir():
            continue
        if target_subject and subject_dir.name.lower() != target_subject.lower():
            continue
        print(f"\n{'#'*80}")
        print(f"Processing subject: {subject_dir.name}")
        print(f"{'#'*80}\n")

        raw_dir = subject_dir / "raw"
        processed_dir = subject_dir / "processed"

        if not raw_dir.exists():
            print(f"Warning: No raw directory found for {subject_dir.name}")
            continue

        extraction_tasks = [
            (raw_dir / "Notes", processed_dir / "notes.txt"),
            (raw_dir / "Question Papers", processed_dir / "questions.txt"),
            (raw_dir / "Syllabus", processed_dir / "syllabus.txt"),
        ]

        total_processed = 0
        total_succeeded = 0
        for source_dir, output_file in extraction_tasks:
            if source_dir.exists():
                print(f"\nExtracting from: {source_dir}")
                processed, succeeded = walk_and_extract(str(source_dir), str(output_file))
                total_processed += processed
                total_succeeded += succeeded
                print(f"  Files processed: {processed}, Successfully extracted: {succeeded}")
            else:
                print(f"Warning: Directory not found: {source_dir}")
        
        print(f"\nSummary for {subject_dir.name}:")
        print(f"  Total files processed: {total_processed}")
        print(f"  Total files succeeded: {total_succeeded}")
    
    print(f"\n{'#'*80}")
    print("Extraction complete!")
    print(f"{'#'*80}")

if __name__ == "__main__":
    missing_libs = []
    if PyPDF2 is None:
        missing_libs.append("PyPDF2")
    if Document is None:
        missing_libs.append("python-docx")
    if Presentation is None:
        missing_libs.append("python-pptx")
    
    if missing_libs:
        print("Warning: Missing libraries for full functionality:")
        for lib in missing_libs:
            print(f"  - {lib}")
        print("\nInstall missing libraries with:")
        print(f"  pip install {' '.join(missing_libs)}")
        print("\nContinuing with available extractors...\n")
    
    main()